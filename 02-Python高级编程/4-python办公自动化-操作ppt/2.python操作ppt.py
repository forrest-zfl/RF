from pptx import Presentation

# 创建幻灯片对象
pres = Presentation()
# 选择母版添加一页
tilte_slide_layout = pres.slide_layouts[0]
slide = pres.slides.add_slide(tilte_slide_layout)

# 获取标题栏和副标题栏
title = slide.shapes.title
subtitle = slide.placeholders[1]

# 编辑主标题和副标题
title.text = "welcome to python"
subtitle.text = "Life is short, I use python"


# 选择母版添加一页
bullet_slide_layout = pres.slide_layouts[1]
slide = pres.slides.add_slide(bullet_slide_layout)

# 获取页面上面所有的形状
shapes = slide.shapes
# print(shapes)

# 获取标题和主体
title_shape = shapes.title
body_shape = shapes.placeholders[1]

# 编辑标题
title_shape.text = "介绍"
# 编辑主体内容
tf = body_shape.text_frame
tf.text = "History of Python"

# 添加一个一级段落
p = tf.add_paragraph()
p.text = "hello world"
p.level = 1

# 添加一个二级段落
p = tf.add_paragraph()
p.text = "马上开始"
p.level = 2

# 保存幻灯片
pres.save("测试.pptx")